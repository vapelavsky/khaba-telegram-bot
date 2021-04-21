from pptx import Presentation
from pptx.util import Inches


def create_presentation(name_list, image_list, faculty):
    prs = Presentation()
    for name in name_list:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = name
    for image in image_list:
        left = Inches(1)
        top = Inches(2)
        width = Inches(3)
        height = Inches(3)
        img = slide.shapes.add_picture(image, left, top, width, height)

    prs.save(f'{faculty}_report.pptx')
