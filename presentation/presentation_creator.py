import os

from pptx import Presentation
from pptx.util import Inches
from db.config import BASE_DIR
from celery import Celery

app = Celery() # Ініціалізуємо Celery


@app.task #Декоратор, який вносить задачу створення презентації до асинхронної черги задач 
def create_presentation(name, image_list, faculty): # Виклик функції, що генерує презентацію
    prs = Presentation()
    steps = [Inches(1), Inches(4), Inches(6.5)] # Кроки для переміщення зображень на слайді
    for i in range(len(image_list)): # Цикл, що ітерується по довжині масиву списку зображень
        slide = prs.slides.add_slide(prs.slide_layouts[1]) 
        title_shape = slide.shapes.title
        title_shape.text = name[i] # Додавання назви заходу на слайд
        for image in image_list[i]: # Цикл, що додає зображення на слайд
            left = steps[image_list[i].index(image)] #На скільки кроків рухається зображення
            # Розміри зображення
            top = Inches(2)
            width = Inches(3)
            height = Inches(3)
            img = slide.shapes.add_picture(image, left, top, width, height) # Додавання зображення згідно заданих параметрів

    prs.save(os.path.join(BASE_DIR, f"presentations/{faculty}_report.pptx")) # Збереження презентації з назвою факультету, для якого генерується звіт


if __name__ == "__main__":
    app.worker_main() # Запуск задачі задачі у Celery 
